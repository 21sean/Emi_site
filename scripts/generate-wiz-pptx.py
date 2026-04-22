#!/usr/bin/env python3
"""Generate Google × Wiz M&A presentation as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(SCRIPT_DIR, "assets")
OUTPUT = os.path.join(SCRIPT_DIR, "..", "public", "google-wiz-ma.pptx")

# Colors
NAVY      = RGBColor(0x0f, 0x17, 0x2a)
DARK_BLUE = RGBColor(0x1e, 0x29, 0x3b)
MID_BLUE  = RGBColor(0x33, 0x41, 0x55)
SLATE     = RGBColor(0x64, 0x74, 0x8b)
ACCENT    = RGBColor(0x25, 0x63, 0xeb)
GOLD      = RGBColor(0xf5, 0x9e, 0x0b)
GOLD_DARK = RGBColor(0xd9, 0x77, 0x06)
RED_SOFT  = RGBColor(0xef, 0x44, 0x44)
GREEN     = RGBColor(0x22, 0xc5, 0x5e)
WHITE     = RGBColor(0xff, 0xff, 0xff)
TEXT_DARK  = RGBColor(0x0f, 0x17, 0x2a)
TEXT_MID   = RGBColor(0x33, 0x41, 0x55)
TEXT_LIGHT = RGBColor(0x64, 0x74, 0x8b)
BORDER     = RGBColor(0xe2, 0xe8, 0xf0)
CARD_BG    = RGBColor(0xf1, 0xf5, 0xf9)
LIGHT_BG   = RGBColor(0xf8, 0xfa, 0xfc)

def asset(name):
    return os.path.join(ASSETS, name)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=TEXT_MID, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=10, color=TEXT_MID, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox

def add_page_number(slide, num, total=22, color=SLATE):
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
                f"{num} / {total}", font_size=8, color=color, alignment=PP_ALIGN.RIGHT)

def add_accent_line(slide, left, top, width, color=ACCENT):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line

def add_top_bar(slide, color=NAVY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False

def add_image_safe(slide, path, left, top, width=None, height=None):
    try:
        slide.shapes.add_picture(path, left, top, width=width, height=height)
    except Exception as e:
        print(f"Warning: could not add image {path}: {e}")


# ═══════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # Gold top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Pt(8))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background(); bar.shadow.inherit = False

    # Google logo
    add_image_safe(slide, asset("google_logo.png"), Inches(0.7), Inches(0.5), width=Inches(2.5))

    # × symbol
    add_textbox(slide, Inches(3.3), Inches(0.5), Inches(0.5), Inches(0.6), "×",
                font_size=32, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Wiz logo
    add_image_safe(slide, asset("wiz_logo.png"), Inches(3.9), Inches(0.5), width=Inches(1.6))

    # Shield + cloud icons (right side)
    add_image_safe(slide, asset("cloud.png"), Inches(9.5), Inches(1.5), width=Inches(1.8))
    add_image_safe(slide, asset("shield.png"), Inches(11), Inches(1.2), width=Inches(1.5))

    # Title
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(8), Inches(0.8),
                "M&A Strategic &", font_size=40, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.7), Inches(2.9), Inches(8), Inches(0.8),
                "Valuation Analysis", font_size=40, bold=True, color=WHITE)

    add_accent_line(slide, Inches(0.7), Inches(3.7), Inches(2), GOLD)

    # $32B callout
    box = add_shape(slide, Inches(0.7), Inches(4.0), Inches(2.5), Inches(0.8),
                    fill_color=DARK_BLUE, line_color=MID_BLUE)
    add_textbox(slide, Inches(0.85), Inches(4.05), Inches(1.2), Inches(0.7),
                "$32B", font_size=28, bold=True, color=GOLD)
    add_textbox(slide, Inches(2.1), Inches(4.15), Inches(1), Inches(0.5),
                "All-Cash", font_size=12, color=SLATE)

    # Merger arrows
    add_image_safe(slide, asset("merger.png"), Inches(3.5), Inches(4.0), width=Inches(1.5))

    # Authors
    add_textbox(slide, Inches(0.7), Inches(5.8), Inches(10), Inches(0.4),
                "Ryan Benirschke  ·  Stone Chung  ·  Alicia Hsiao  ·  Emi Kobayashi  ·  Dean Wu",
                font_size=13, color=SLATE)

    add_textbox(slide, Inches(0.7), Inches(6.2), Inches(8), Inches(0.3),
                "Mergers & Acquisitions  |  Fall 2025  |  UC San Diego",
                font_size=11, color=SLATE)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(5), Inches(0.6),
                "Agenda", font_size=32, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    items = [
        ("01", "Background", "Industry landscape & competitive positioning", "cloud.png"),
        ("02", "Google's Acquisition", "Timeline, synergies & strategic rationale", "merger.png"),
        ("03", "Strategic Outlook", "Market gaps, consolidation & alternatives", "chart.png"),
        ("04", "Acquisition Deal Structure", "Valuation, premiums & regulatory review", "money.png"),
        ("05", "Our Team Analysis", "Final verdict & conclusions", "shield.png"),
    ]

    for i, (num, title, desc, icon_name) in enumerate(items):
        y = Inches(1.4) + i * Inches(1.1)

        # Number circle
        circ = add_shape(slide, Inches(0.7), y, Inches(0.55), Inches(0.55),
                         fill_color=ACCENT if i == 0 else CARD_BG, shape_type=MSO_SHAPE.OVAL)
        add_textbox(slide, Inches(0.7), y + Pt(5), Inches(0.55), Inches(0.45),
                    num, font_size=14, bold=True, color=WHITE if i == 0 else ACCENT, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.5), y, Inches(8), Inches(0.35),
                    title, font_size=17, bold=True, color=TEXT_DARK)
        add_textbox(slide, Inches(1.5), y + Inches(0.33), Inches(8), Inches(0.3),
                    desc, font_size=11, color=TEXT_LIGHT)

        # Icon
        add_image_safe(slide, asset(icon_name), Inches(11.5), y, width=Inches(0.6))

        # Divider
        if i < len(items) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(1.5), y + Inches(0.8), Inches(10.5), Pt(0.5))
            line.fill.solid(); line.fill.fore_color.rgb = BORDER
            line.line.fill.background(); line.shadow.inherit = False

    add_page_number(slide, 2)


def slide_section(prs, number, title, page_num, icon_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Gold top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background(); bar.shadow.inherit = False

    # Large number watermark
    add_textbox(slide, Inches(0.3), Inches(2.2), Inches(3.5), Inches(2),
                f"0{number}", font_size=120, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)

    # Title
    add_textbox(slide, Inches(3), Inches(3.0), Inches(7), Inches(0.8),
                title, font_size=36, bold=True, color=WHITE)
    add_accent_line(slide, Inches(3), Inches(3.75), Inches(1.5), GOLD)

    if icon_name:
        add_image_safe(slide, asset(icon_name), Inches(10.5), Inches(2.7), width=Inches(1.5))

    add_page_number(slide, page_num, color=MID_BLUE)


def slide_industry(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
                "Industry Background", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    add_image_safe(slide, asset("cloud.png"), Inches(11.3), Inches(0.3), width=Inches(1.2))

    add_textbox(slide, Inches(0.7), Inches(1.2), Inches(11), Inches(0.6),
                "Wiz is a cloud security platform that helps organizations identify, prioritize, and remediate risks "
                "across their cloud environments through agentless, multi-cloud visibility and analysis.",
                font_size=11, color=TEXT_MID)

    # CSPM card
    cspm = add_shape(slide, Inches(0.7), Inches(2.0), Inches(5.6), Inches(1.1),
                     fill_color=RGBColor(0xef, 0xf6, 0xff), line_color=RGBColor(0xbf, 0xdb, 0xfe))
    add_image_safe(slide, asset("shield.png"), Inches(5.4), Inches(2.05), width=Inches(0.6))
    add_textbox(slide, Inches(0.9), Inches(2.1), Inches(4.5), Inches(0.3),
                "Cloud Security Posture Management (CSPM)", font_size=11, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.9), Inches(2.45), Inches(4.5), Inches(0.5),
                "Continuous monitoring of cloud infrastructure for misconfigurations and compliance violations",
                font_size=9, color=TEXT_MID)

    # CNAPP card
    cnapp = add_shape(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(1.1),
                      fill_color=RGBColor(0xfe, 0xf3, 0xc7), line_color=RGBColor(0xfd, 0xe6, 0x8a))
    add_image_safe(slide, asset("cloud.png"), Inches(11.5), Inches(2.1), width=Inches(0.6))
    add_textbox(slide, Inches(7.0), Inches(2.1), Inches(4.5), Inches(0.3),
                "Cloud-Native Application Protection (CNAPP)", font_size=11, bold=True, color=GOLD_DARK)
    add_textbox(slide, Inches(7.0), Inches(2.45), Inches(4.5), Inches(0.5),
                "End-to-end security for cloud-native applications from development through runtime",
                font_size=9, color=TEXT_MID)

    # Porter's Five Forces
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(5), Inches(0.4),
                "Porter's Five Forces Analysis", font_size=14, bold=True, color=TEXT_DARK)

    forces = [
        ("Threat of New\nEntrants", "Moderate", GOLD),
        ("Threat of\nSubstitutes", "Moderate", GOLD),
        ("Competitive\nRivalry", "High", RED_SOFT),
        ("Buyer Bargaining\nPower", "Moderate", GOLD),
        ("Supplier Bargaining\nPower", "Low", GREEN),
    ]

    for i, (label, level, color) in enumerate(forces):
        x = Inches(0.7) + i * Inches(2.4)
        y = Inches(4.1)

        card = add_shape(slide, x, y, Inches(2.1), Inches(1.4), fill_color=LIGHT_BG, line_color=BORDER)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.9), Inches(0.6),
                    label, font_size=9, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

        badge = add_shape(slide, x + Inches(0.5), y + Inches(0.9), Inches(1.1), Inches(0.35),
                          fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, x + Inches(0.5), y + Inches(0.92), Inches(1.1), Inches(0.3),
                    level, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 4)


def slide_strategic_positioning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.6),
                "Strategic Positioning & Acquisition Rationale", font_size=24, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    cols = [
        ("Competitive Positioning", RGBColor(0xef, 0xf6, 0xff), RGBColor(0xbf, 0xdb, 0xfe), ACCENT, "shield.png",
         ["Agentless architecture", "Fast, frictionless deployment", "Full cloud-wide visibility", "50% Fortune 100 penetration"]),
        ("Why Google", RGBColor(0xf0, 0xfd, 0xf4), RGBColor(0xbb, 0xf7, 0xd0), GREEN, "google_logo.png",
         ["Accelerate scaling through Google Cloud ecosystem", "Expand R&D using Google's AI and security infrastructure", "Enable high-value liquidity event for investors"]),
        ("Wiz's Alternatives", RGBColor(0xfe, 0xf3, 0xc7), RGBColor(0xfd, 0xe6, 0x8a), GOLD_DARK, "wiz_logo.png",
         ["Remain independent", "Pursue an IPO", "Form strategic alliances", "Engage other acquirers (AWS, Microsoft, Palo Alto, CrowdStrike)"]),
    ]

    for i, (title, bg, border, accent, icon, bullets) in enumerate(cols):
        x = Inches(0.7) + i * Inches(4.1)
        y = Inches(1.3)

        card = add_shape(slide, x, y, Inches(3.8), Inches(5.5), fill_color=bg, line_color=border)
        add_image_safe(slide, asset(icon), x + Inches(1.4), y + Inches(0.2), width=Inches(1))
        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.4), Inches(0.4),
                    title, font_size=13, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.5), Inches(3.2), Inches(3.5),
                        bullets, font_size=10, color=TEXT_MID)

    add_page_number(slide, 5)


def slide_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
                "Acquisition Timeline", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)
    add_image_safe(slide, asset("merger.png"), Inches(11.3), Inches(0.3), width=Inches(1.3))

    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.0), Inches(11), Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background(); line.shadow.inherit = False

    events = [
        ("2020", "Wiz Founded", "Cloud security startup launched by former Microsoft engineers", False),
        ("2021–2022", "Hypergrowth", "ARR surges from $1M to $100M in just 18 months", False),
        ("2024", "Google's First Offer", "Wiz rejects Google Cloud's $23B acquisition offer", False),
        ("2025", "$32B Acquisition", "Google and Wiz announce all-cash agreement — Google's largest acquisition ever", True),
    ]

    for i, (year, title, desc, highlight) in enumerate(events):
        x = Inches(1.2) + i * Inches(3.0)

        # Dot on timeline
        dot = add_shape(slide, x + Inches(0.8), Inches(2.85), Inches(0.3), Inches(0.3),
                        fill_color=GOLD if highlight else ACCENT, shape_type=MSO_SHAPE.OVAL)

        # Year above
        add_textbox(slide, x + Inches(0.3), Inches(2.3), Inches(1.4), Inches(0.4),
                    year, font_size=14, bold=True, color=GOLD if highlight else ACCENT, alignment=PP_ALIGN.CENTER)

        # Card below
        bg = RGBColor(0xfe, 0xf3, 0xc7) if highlight else CARD_BG
        brd = RGBColor(0xfd, 0xe6, 0x8a) if highlight else BORDER
        card = add_shape(slide, x, Inches(3.5), Inches(2.5), Inches(2.0), fill_color=bg, line_color=brd)

        add_textbox(slide, x + Inches(0.15), Inches(3.6), Inches(2.2), Inches(0.35),
                    title, font_size=11, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(4.0), Inches(2.2), Inches(1.3),
                    desc, font_size=9, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    add_image_safe(slide, asset("chart.png"), Inches(11.5), Inches(6.2), width=Inches(1))
    add_page_number(slide, 7)


def slide_synergies(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
                "Acquisition Synergies", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    synergies = [
        ("Innovation Acceleration",
         "Integration into Google Cloud accelerates product development through AI infrastructure and resources.",
         "→ Faster innovation and time-to-market",
         RGBColor(0xef, 0xf6, 0xff), ACCENT, "chart.png"),
        ("AI-Driven Cloud Security",
         "Investment focuses on advancing AI-era cloud security and multi-cloud capabilities.",
         "→ Improved security design, automation, and threat prevention across AI workloads",
         RGBColor(0xf0, 0xfd, 0xf4), GREEN, "shield.png"),
        ("Operational Efficiency",
         "Automated platform scales security operations and reduces manual workload.",
         "→ Lower total cost of ownership and improved management efficiency",
         RGBColor(0xfe, 0xf3, 0xc7), GOLD_DARK, "cloud.png"),
        ("Market Expansion",
         "Wiz's enterprise and startup customer base complements Google's market reach.",
         "→ Broader adoption of Google Cloud and enhanced enterprise growth",
         RGBColor(0xfa, 0xf5, 0xff), RGBColor(0xa8, 0x55, 0xf7), "merger.png"),
    ]

    for i, (title, desc, benefit, bg, accent, icon) in enumerate(synergies):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + col * Inches(6.2)
        y = Inches(1.3) + row * Inches(2.8)

        card = add_shape(slide, x, y, Inches(5.8), Inches(2.5), fill_color=bg, line_color=BORDER)

        # Accent left bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(5), Inches(2.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False

        add_image_safe(slide, asset(icon), x + Inches(4.8), y + Inches(0.2), width=Inches(0.6))

        add_textbox(slide, x + Inches(0.3), y + Inches(0.15), Inches(4), Inches(0.35),
                    title, font_size=13, bold=True, color=accent)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.55), Inches(4.5), Inches(1.0),
                    desc, font_size=10, color=TEXT_MID)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.9), Inches(5), Inches(0.3),
                    benefit, font_size=9, bold=True, color=accent)

    add_page_number(slide, 8)


def slide_strat_problem_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
                "The Strategic Problem Google Must Solve", font_size=22, bold=True, color=TEXT_DARK)
    add_textbox(slide, Inches(0.7), Inches(0.8), Inches(10), Inches(0.4),
                "Google's Structural Gap in Multi-Cloud Security", font_size=15, bold=True, color=ACCENT)
    add_accent_line(slide, Inches(0.7), Inches(1.15), Inches(1.2), ACCENT)
    add_image_safe(slide, asset("cloud.png"), Inches(11.3), Inches(0.3), width=Inches(1.2))

    # Warning box
    warn = add_shape(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.9),
                     fill_color=RGBColor(0xfe, 0xf2, 0xf2), line_color=RGBColor(0xfe, 0xca, 0xca))
    add_textbox(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(0.7),
                "Customers increasingly demand unified security visibility across all clouds, "
                "while Google's security tools remain largely GCP-centric. This places Google at "
                "a structural disadvantage in enterprise cloud adoption.",
                font_size=11, bold=True, color=RGBColor(0x99, 0x1b, 0x1b))

    add_textbox(slide, Inches(0.7), Inches(2.6), Inches(5), Inches(0.4),
                "Why enterprises adopt multi-cloud:", font_size=14, bold=True, color=TEXT_DARK)

    reasons = [
        ("1", "Avoid Vendor Lock-In", "Relying on a single cloud provider creates strategic and pricing risks. Multi-cloud gives enterprises stronger negotiation power and flexibility.", "shield.png"),
        ("2", "Optimize Workloads", "Different clouds excel at different things — AWS for scale, Azure for enterprise integration, GCP for data/AI — so companies place each workload where it performs best.", "cloud.png"),
        ("3", "Increase Resilience", "Using multiple cloud providers prevents a single point of failure and improves business continuity.", "chart.png"),
    ]

    for i, (num, title, desc, icon) in enumerate(reasons):
        x = Inches(0.7) + i * Inches(4.1)
        y = Inches(3.2)

        card = add_shape(slide, x, y, Inches(3.8), Inches(3.5), fill_color=CARD_BG, line_color=BORDER)

        circ = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4),
                         fill_color=ACCENT, shape_type=MSO_SHAPE.OVAL)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.18), Inches(0.4), Inches(0.35),
                    num, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.65), y + Inches(0.2), Inches(2.5), Inches(0.3),
                    title, font_size=12, bold=True, color=TEXT_DARK)

        add_image_safe(slide, asset(icon), x + Inches(3.1), y + Inches(0.1), width=Inches(0.45))

        add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(2.5),
                    desc, font_size=9.5, color=TEXT_MID)

    add_page_number(slide, 10)


def slide_strat_problem_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
                "The Strategic Problem Google Must Solve", font_size=22, bold=True, color=TEXT_DARK)
    add_textbox(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.4),
                "Market Consolidation Accelerates Toward End-to-End Platforms", font_size=15, bold=True, color=ACCENT)
    add_accent_line(slide, Inches(0.7), Inches(1.15), Inches(1.2), ACCENT)

    # Left insights
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5), Inches(0.35),
                "Google (Major Players Quadrant)", font_size=12, bold=True, color=TEXT_DARK)
    add_image_safe(slide, asset("google_logo.png"), Inches(4.2), Inches(1.5), width=Inches(0.9))

    add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(6), Inches(2),
                    ["Strong on strategy but lacking in product depth",
                     "Multi-cloud visibility and integration remain limited",
                     "Enterprises do not view Google as a first-choice security platform",
                     "Lacks a fully integrated, end-to-end CNAPP solution"],
                    font_size=10, color=TEXT_MID)

    add_textbox(slide, Inches(0.7), Inches(3.8), Inches(5), Inches(0.35),
                "Wiz (Leaders Quadrant)", font_size=12, bold=True, color=TEXT_DARK)
    add_image_safe(slide, asset("wiz_logo.png"), Inches(3.2), Inches(3.8), width=Inches(0.7))

    add_bullet_list(slide, Inches(0.7), Inches(4.2), Inches(6), Inches(2),
                    ["Most complete multi-cloud coverage",
                     "Fastest enterprise adoption (ARR growing over 100%)",
                     "Leading AI-driven and graph-based cloud security",
                     "Best-in-class technology capability and strategic execution"],
                    font_size=10, color=TEXT_MID)

    # Quadrant diagram (right side)
    quad = add_shape(slide, Inches(7.5), Inches(1.4), Inches(5.2), Inches(5.2),
                     fill_color=LIGHT_BG, line_color=BORDER)

    # Axis labels
    add_textbox(slide, Inches(9.5), Inches(6.35), Inches(1.5), Inches(0.25),
                "Strategies →", font_size=8, color=SLATE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.6), Inches(1.5), Inches(0.5), Inches(0.25),
                "Capabilities ↑", font_size=8, color=SLATE)

    # Quadrant labels
    add_textbox(slide, Inches(7.8), Inches(1.7), Inches(1.5), Inches(0.25),
                "Niche Players", font_size=8, color=TEXT_LIGHT)
    add_textbox(slide, Inches(11.3), Inches(1.7), Inches(1.2), Inches(0.25),
                "Leaders", font_size=8, color=TEXT_LIGHT, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(7.8), Inches(6.1), Inches(1.2), Inches(0.25),
                "Emerging", font_size=8, color=TEXT_LIGHT)
    add_textbox(slide, Inches(11.3), Inches(6.1), Inches(1.2), Inches(0.25),
                "Major Players", font_size=8, color=TEXT_LIGHT, alignment=PP_ALIGN.RIGHT)

    # Company dots
    wiz_dot = add_shape(slide, Inches(11.4), Inches(2.0), Inches(0.25), Inches(0.25),
                        fill_color=GOLD, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(11.7), Inches(2.0), Inches(0.6), Inches(0.25),
                "Wiz", font_size=8, bold=True, color=GOLD_DARK)

    pa_dot = add_shape(slide, Inches(10.8), Inches(2.4), Inches(0.2), Inches(0.2),
                       fill_color=RGBColor(0xa8, 0x55, 0xf7), shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(11.05), Inches(2.4), Inches(0.8), Inches(0.2),
                "Palo Alto", font_size=7, bold=True, color=RGBColor(0x7c, 0x3a, 0xed))

    g_dot = add_shape(slide, Inches(10.2), Inches(4.0), Inches(0.25), Inches(0.25),
                      fill_color=ACCENT, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(10.5), Inches(4.0), Inches(0.8), Inches(0.25),
                "Google", font_size=8, bold=True, color=ACCENT)

    add_page_number(slide, 11)


def slide_alternatives_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
                "Evaluating Google's Strategic Alternatives", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    options = [
        ("Internal Development", RGBColor(0xef, 0xf6, 0xff), RGBColor(0xbf, 0xdb, 0xfe), ACCENT, "chart.png", False,
         ["Full control over design and roadmap", "Slow: 3–5+ years to reach parity with Wiz",
          "Risk of falling further behind in a fast-moving market", "No access to Wiz's existing customer relationships"]),
        ("Partnership", RGBColor(0xf0, 0xfd, 0xf4), RGBColor(0xbb, 0xf7, 0xd0), GREEN, "merger.png", False,
         ["Lower cost and lower risk than acquisition", "Limited integration depth and control",
          "Partner may later align with a competitor", "Does not close the structural multi-cloud gap"]),
        ("Acquisition", RGBColor(0xfe, 0xf3, 0xc7), RGBColor(0xfd, 0xe6, 0x8a), GOLD_DARK, "money.png", True,
         ["Immediate access to Wiz's technology and customers", "Closes the multi-cloud security gap overnight",
          "High upfront cost ($32B all-cash)", "Integration and retention risks"]),
    ]

    for i, (title, bg, border, accent, icon, selected, bullets) in enumerate(options):
        x = Inches(0.7) + i * Inches(4.1)
        y = Inches(1.3)

        card = add_shape(slide, x, y, Inches(3.8), Inches(5.5),
                         fill_color=bg, line_color=accent if selected else border)

        if selected:
            badge = add_shape(slide, x + Inches(2.6), y + Inches(0.1), Inches(1.1), Inches(0.35),
                              fill_color=accent, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, x + Inches(2.6), y + Inches(0.12), Inches(1.1), Inches(0.3),
                        "SELECTED", font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_image_safe(slide, asset(icon), x + Inches(1.4), y + Inches(0.3), width=Inches(0.8))
        add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.4), Inches(0.4),
                    title, font_size=14, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.7), Inches(3.2), Inches(3.5),
                        bullets, font_size=10, color=TEXT_MID)

    add_page_number(slide, 12)


def slide_alternatives_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
                "Strategic Alternatives Comparison", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    rows_data = [
        ["Criteria", "Internal Dev", "Partnership", "Acquisition"],
        ["Time to market", "3–5+ years", "1–2 years", "Immediate"],
        ["Cost", "Low–Medium", "Low", "Very High ($32B)"],
        ["Control", "Full", "Limited", "Full"],
        ["Multi-cloud gap", "Gradual close", "Partial", "Fully closed"],
        ["Customer access", "None", "Shared", "Full (50% F100)"],
        ["Risk level", "Execution risk", "Alignment risk", "Integration risk"],
        ["Strategic impact", "Low–Medium", "Medium", "Transformative"],
    ]

    table_shape = slide.shapes.add_table(len(rows_data), 4, Inches(1.5), Inches(1.4), Inches(10), Inches(5))
    table = table_shape.table

    for col_idx in range(4):
        table.columns[col_idx].width = Inches(2.5)

    for row_idx, row in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10 if row_idx > 0 else 11)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY if col_idx < 3 else ACCENT
            elif col_idx == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT
            elif col_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_DARK
            else:
                p.font.color.rgb = TEXT_MID

            if row_idx > 0 and row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    add_page_number(slide, 13)


def slide_deal_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.6),
                "Google–Wiz Acquisition Overview & Deal Structure", font_size=24, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    # Left card
    left = add_shape(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5), fill_color=LIGHT_BG, line_color=BORDER)
    add_textbox(slide, Inches(0.9), Inches(1.4), Inches(2), Inches(0.4),
                "Deal Overview", font_size=15, bold=True, color=ACCENT)
    add_image_safe(slide, asset("google_logo.png"), Inches(4.8), Inches(1.35), width=Inches(1.3))

    details = [
        ("Target", "Wiz — cloud security platform (founded 2020)"),
        ("Acquirer", "Google Cloud (Alphabet)"),
        ("Announced", "March 18, 2025"),
        ("Deal Value", "US $32B — Google's largest acquisition ever"),
        ("Expected Close", "Pending regulatory approvals; expected 2026"),
    ]
    for i, (label, value) in enumerate(details):
        y = Inches(2.0) + i * Inches(0.4)
        add_textbox(slide, Inches(0.9), y, Inches(1.3), Inches(0.3), label, font_size=10, bold=True, color=TEXT_DARK)
        add_textbox(slide, Inches(2.3), y, Inches(4), Inches(0.3), value, font_size=10, color=TEXT_MID)

    add_image_safe(slide, asset("money.png"), Inches(2.8), Inches(4.8), width=Inches(1))

    # Right card
    right = add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                      fill_color=RGBColor(0xfe, 0xf3, 0xc7), line_color=RGBColor(0xfd, 0xe6, 0x8a))
    add_textbox(slide, Inches(7.0), Inches(1.4), Inches(3.5), Inches(0.4),
                "Deal Financing & Structure", font_size=15, bold=True, color=GOLD_DARK)
    add_image_safe(slide, asset("wiz_logo.png"), Inches(11.3), Inches(1.35), width=Inches(0.9))

    add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.2), Inches(3.5),
                    ["Deal nature: Friendly acquisition; negotiated and approved by both boards",
                     "Financing: All-cash transaction (no stock, no share swap)",
                     "Context: Google previously offered $23B in 2024, which Wiz rejected; final agreement reached at $32B in 2025",
                     "Integration: Wiz leadership to stay on and guide integration under Google Cloud"],
                    font_size=10, color=TEXT_MID)

    add_image_safe(slide, asset("merger.png"), Inches(9.0), Inches(5.2), width=Inches(1.2))

    add_textbox(slide, Inches(0.7), Inches(6.9), Inches(5), Inches(0.25),
                "Source: Google, SEC, Nasdaq", font_size=8, color=TEXT_LIGHT)
    add_page_number(slide, 15)


def slide_premium(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
                "Premium & Revenue Multiple Analysis", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)
    add_image_safe(slide, asset("money.png"), Inches(11.5), Inches(0.2), width=Inches(0.9))

    # Left: Premium cards
    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(5), Inches(0.4),
                "Premium Over Prior Valuations", font_size=14, bold=True, color=TEXT_DARK)

    premiums = [
        ("167%", "over $12B May 2024 funding round", RGBColor(0xfe, 0xf2, 0xf2), RED_SOFT),
        ("100%", "over $16B late-2024 secondary market valuation", RGBColor(0xfe, 0xf3, 0xc7), GOLD_DARK),
        ("39%", "over Google's first $23B offer (mid-2024)", RGBColor(0xef, 0xf6, 0xff), ACCENT),
    ]

    for i, (pct, desc, bg, color) in enumerate(premiums):
        y = Inches(1.8) + i * Inches(1.1)
        card = add_shape(slide, Inches(0.7), y, Inches(5.8), Inches(0.85), fill_color=bg, line_color=BORDER)
        add_textbox(slide, Inches(0.9), y + Inches(0.15), Inches(1.2), Inches(0.5),
                    pct, font_size=24, bold=True, color=color)
        add_textbox(slide, Inches(2.2), y + Inches(0.25), Inches(4), Inches(0.4),
                    desc, font_size=10, color=TEXT_MID)

    # Right: Revenue multiples
    add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
                "Implied Revenue Multiples at $32B", font_size=14, bold=True, color=TEXT_DARK)

    big_box = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(1.1), fill_color=NAVY)
    add_textbox(slide, Inches(7.0), Inches(1.85), Inches(1.5), Inches(0.7),
                "32×", font_size=36, bold=True, color=GOLD)
    add_textbox(slide, Inches(8.5), Inches(1.95), Inches(3.5), Inches(0.35),
                "forward revenue", font_size=13, color=WHITE)
    add_textbox(slide, Inches(8.5), Inches(2.35), Inches(3.5), Inches(0.3),
                "Based on projected $1B ARR in 2025", font_size=10, color=SLATE)

    # Comparison bars
    comparisons = [
        ("Typical cybersecurity deals", "5–15×", 0.25, ACCENT),
        ("Auth0 (prior high watermark)", "42×", 0.85, ACCENT),
        ("Google–Wiz", "32×", 0.65, GOLD),
    ]

    for i, (label, mult, pct, color) in enumerate(comparisons):
        y = Inches(3.3) + i * Inches(0.9)
        add_textbox(slide, Inches(6.8), y, Inches(5), Inches(0.25), label, font_size=10, color=TEXT_MID)

        bar_bg = add_shape(slide, Inches(6.8), y + Inches(0.3), Inches(5.5), Inches(0.3),
                           fill_color=CARD_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        bar_fill = add_shape(slide, Inches(6.8), y + Inches(0.3), Inches(5.5 * pct), Inches(0.3),
                             fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, Inches(6.8) + Inches(5.5 * pct) - Inches(0.6), y + Inches(0.3), Inches(0.6), Inches(0.3),
                    mult, font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Note
    note = add_shape(slide, Inches(6.8), Inches(6.0), Inches(5.8), Inches(0.6),
                     fill_color=RGBColor(0xfe, 0xf3, 0xc7), line_color=RGBColor(0xfd, 0xe6, 0x8a))
    add_textbox(slide, Inches(7.0), Inches(6.05), Inches(5.4), Inches(0.5),
                "Google–Wiz is clearly priced as a future multi-billion ARR platform, not a typical security vendor.",
                font_size=9, color=GOLD_DARK)

    add_page_number(slide, 16)


def slide_regulatory(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
                "Regulatory & Legal Structure", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    items = [
        ("01", "Regulatory Review",
         "United States Department of Justice (DOJ) cleared antitrust review for the deal.",
         RGBColor(0xef, 0xf6, 0xff), ACCENT, "shield.png"),
        ("02", "Integration Commitment",
         "Google states that Wiz's products will continue to work across all major clouds (multi-cloud support) to ease regulatory concerns.",
         RGBColor(0xf0, 0xfd, 0xf4), GREEN, "cloud.png"),
        ("03", "Closing Conditions",
         "Subject to regulatory approvals around the world. Expected to close in 2026.",
         RGBColor(0xfe, 0xf3, 0xc7), GOLD_DARK, "money.png"),
    ]

    for i, (num, title, desc, bg, accent, icon) in enumerate(items):
        x = Inches(0.7) + i * Inches(4.1)
        y = Inches(1.3)

        card = add_shape(slide, x, y, Inches(3.8), Inches(5.5), fill_color=bg, line_color=BORDER)

        add_image_safe(slide, asset(icon), x + Inches(1.5), y + Inches(0.3), width=Inches(0.7))

        circ = add_shape(slide, x + Inches(1.55), y + Inches(1.3), Inches(0.6), Inches(0.6),
                         fill_color=accent, shape_type=MSO_SHAPE.OVAL)
        add_textbox(slide, x + Inches(1.55), y + Inches(1.35), Inches(0.6), Inches(0.5),
                    num, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(2.1), Inches(3.4), Inches(0.4),
                    title, font_size=14, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y + Inches(2.7), Inches(3.2), Inches(2.5),
                    desc, font_size=10.5, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.7), Inches(6.9), Inches(5), Inches(0.25),
                "Source: Yahoo Finance", font_size=8, color=TEXT_LIGHT)
    add_page_number(slide, 17)


def slide_good_deal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
                "Was This a Good Deal?", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)

    # Financial (left)
    left = add_shape(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5),
                     fill_color=RGBColor(0xfe, 0xf2, 0xf2), line_color=RGBColor(0xfe, 0xca, 0xca))
    add_image_safe(slide, asset("money.png"), Inches(3.0), Inches(1.5), width=Inches(0.8))
    add_textbox(slide, Inches(0.9), Inches(2.3), Inches(5.4), Inches(0.4),
                "Financially — Overpriced", font_size=16, bold=True, color=RED_SOFT, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, Inches(0.9), Inches(2.9), Inches(5.4), Inches(3.5),
                    ["Our team believes this price was too expensive",
                     "The initial $23B offer would have been a much better deal",
                     "Overall short-term investor sentiment was negative — Google's stock dipped as much as 5% at announcement",
                     "This acquisition cost more than Google's previous 8 acquisitions combined"],
                    font_size=10.5, color=TEXT_MID)

    # Strategic (right)
    right = add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                      fill_color=RGBColor(0xf0, 0xfd, 0xf4), line_color=RGBColor(0xbb, 0xf7, 0xd0))
    add_image_safe(slide, asset("shield.png"), Inches(9.2), Inches(1.4), width=Inches(0.7))
    add_textbox(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(0.4),
                "Strategically — Necessary", font_size=16, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, Inches(7.0), Inches(2.9), Inches(5.4), Inches(3.5),
                    ["Critical acquisition given Wiz's existing market penetration and Fortune 500 client base",
                     "Allows Google to become a key player in cloud security",
                     "Provides access to high-value enterprise clients on AWS and Azure",
                     "No antitrust violations as of November 2025"],
                    font_size=10.5, color=TEXT_MID)

    add_page_number(slide, 19)


def slide_final_decision(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background(); bar.shadow.inherit = False

    # Logos
    add_image_safe(slide, asset("google_logo.png"), Inches(4.5), Inches(0.3), width=Inches(2))
    add_textbox(slide, Inches(6.6), Inches(0.35), Inches(0.4), Inches(0.5),
                "×", font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, asset("wiz_logo.png"), Inches(7.1), Inches(0.35), width=Inches(1.3))

    add_textbox(slide, Inches(2), Inches(1.0), Inches(9), Inches(0.7),
                "Final Decision", font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    badge = add_shape(slide, Inches(5.6), Inches(1.7), Inches(2.2), Inches(0.45),
                      fill_color=GOLD, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(5.6), Inches(1.73), Inches(2.2), Inches(0.4),
                "ACQUISITION", font_size=13, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    points = [
        ("Multi-Cloud Access",
         "Google gains immediate access to AWS and Azure workloads through Wiz's native multi-cloud design.",
         "cloud.png"),
        ("AI Capabilities",
         "Wiz benefits from Google's AI capabilities, threat intelligence, and global distribution.",
         "shield.png"),
        ("Industry Leadership",
         "Together, they can lead the industry shift toward unified cloud security platforms.",
         "chart.png"),
    ]

    for i, (title, desc, icon) in enumerate(points):
        x = Inches(0.7) + i * Inches(4.1)
        y = Inches(2.5)

        card = add_shape(slide, x, y, Inches(3.8), Inches(4.0), fill_color=DARK_BLUE, line_color=MID_BLUE)

        # Gold accent line
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.5), y + Inches(0.1), Inches(2.8), Pt(3))
        accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = GOLD; accent_bar.line.fill.background(); accent_bar.shadow.inherit = False

        add_image_safe(slide, asset(icon), x + Inches(1.5), y + Inches(0.4), width=Inches(0.7))

        add_textbox(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.4), Inches(0.4),
                    title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.2), Inches(2.0),
                    desc, font_size=11, color=SLATE, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 20, color=MID_BLUE)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)

    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.6),
                "Conclusion and Outcomes", font_size=26, bold=True, color=TEXT_DARK)
    add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(1.2), ACCENT)
    add_image_safe(slide, asset("merger.png"), Inches(11.3), Inches(0.3), width=Inches(1.3))

    add_bullet_list(slide, Inches(0.7), Inches(1.5), Inches(11), Inches(4.5),
                    ["Google needed to address key gaps in its cloud security business, and Wiz was the ideal target to fill those gaps.",
                     "Wiz has the opportunity to scale rapidly using Google's massive cloud infrastructure and AI capabilities.",
                     "If Google's end goal is to build a market-leading security platform, this acquisition could be the stepping stone they need.",
                     "In general, it is difficult to put a price tag on synergies — short-term investors primarily see the large price tag.",
                     "Given the recency of the acquisition, it remains to be seen how effective this deal will be for Google's strategic position."],
                    font_size=12, color=TEXT_MID)

    # Bottom logos
    add_image_safe(slide, asset("google_logo.png"), Inches(5.2), Inches(6.2), width=Inches(1.3))
    add_textbox(slide, Inches(6.6), Inches(6.25), Inches(0.4), Inches(0.4),
                "×", font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, asset("wiz_logo.png"), Inches(7.1), Inches(6.25), width=Inches(0.9))

    add_page_number(slide, 21)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background(); bar.shadow.inherit = False

    # Logos
    add_image_safe(slide, asset("google_logo.png"), Inches(4.3), Inches(1.5), width=Inches(2.2))
    add_textbox(slide, Inches(6.6), Inches(1.55), Inches(0.5), Inches(0.5),
                "×", font_size=26, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, asset("wiz_logo.png"), Inches(7.2), Inches(1.55), width=Inches(1.5))

    add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(1),
                "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.7), Inches(3.4), Inches(1.8), GOLD)

    add_textbox(slide, Inches(2), Inches(3.7), Inches(9), Inches(0.5),
                "Google × Wiz — M&A Strategic & Valuation Analysis",
                font_size=14, color=SLATE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.4),
                "Ryan Benirschke  ·  Stone Chung  ·  Alicia Hsiao  ·  Emi Kobayashi  ·  Dean Wu",
                font_size=11, color=SLATE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.4),
                "Mergers & Acquisitions  |  Fall 2025  |  UC San Diego",
                font_size=10, color=MID_BLUE, alignment=PP_ALIGN.CENTER)

    add_image_safe(slide, asset("shield.png"), Inches(5.9), Inches(5.5), width=Inches(0.9))


# ═══════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)                                           # 1
    slide_agenda(prs)                                          # 2
    slide_section(prs, 1, "Background", 3, "cloud.png")        # 3
    slide_industry(prs)                                        # 4
    slide_strategic_positioning(prs)                            # 5
    slide_section(prs, 2, "Google's Acquisition", 6, "merger.png")  # 6
    slide_timeline(prs)                                        # 7
    slide_synergies(prs)                                       # 8
    slide_section(prs, 3, "Strategic Outlook", 9, "chart.png") # 9
    slide_strat_problem_1(prs)                                 # 10
    slide_strat_problem_2(prs)                                 # 11
    slide_alternatives_overview(prs)                            # 12
    slide_alternatives_table(prs)                               # 13
    slide_section(prs, 4, "Deal Structure", 14, "money.png")   # 14
    slide_deal_overview(prs)                                   # 15
    slide_premium(prs)                                         # 16
    slide_regulatory(prs)                                      # 17
    slide_section(prs, 5, "Conclusion", 18, "shield.png")      # 18
    slide_good_deal(prs)                                       # 19
    slide_final_decision(prs)                                  # 20
    slide_conclusion(prs)                                      # 21
    slide_thank_you(prs)                                       # 22

    output = os.path.abspath(OUTPUT)
    prs.save(output)
    print(f"Generated: {output}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    build()
